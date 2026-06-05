#!/usr/bin/env python3
"""
PPT Generator v2 - Create professional PowerPoint presentations with
backgrounds, icons, charts, tables, and rich visual elements.

Usage:
    python generate_ppt.py <input.json> <output.pptx> [--theme <theme>]

Themes: blue (default), dark, green, orange, purple, red, minimal
"""

import json
import sys
import os
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.dml.color import RGBColor
    from pptx.chart.data import CategoryChartData
except ImportError:
    print("Error: python-pptx is required. Install with: pip install python-pptx")
    sys.exit(1)

# ── Themes ──────────────────────────────────────────────────────────────

THEMES = {
    "blue": {
        "primary":    RGBColor(0x1A, 0x73, 0xE8),  # Google Blue
        "secondary":  RGBColor(0x0D, 0x47, 0xA1),  # Dark Blue
        "accent":     RGBColor(0x00, 0xBC, 0xD4),  # Cyan
        "bg":         RGBColor(0xFF, 0xFF, 0xFF),
        "bg_dark":    RGBColor(0x0D, 0x47, 0xA1),
        "text":       RGBColor(0x21, 0x21, 0x21),
        "text_light": RGBColor(0xFF, 0xFF, 0xFF),
        "text_muted": RGBColor(0x75, 0x75, 0x75),
        "card_bg":    RGBColor(0xF5, 0xF5, 0xF5),
        "gradient_start": RGBColor(0x1A, 0x73, 0xE8),
        "gradient_end":   RGBColor(0x0D, 0x47, 0xA1),
        "decorations": ["circles", "waves", "dots"],  # 主题装饰风格
        "deco_color1": RGBColor(0x1A, 0x73, 0xE8),
        "deco_color2": RGBColor(0x00, 0xBC, 0xD4),
        "deco_opacity": 0.15,
    },
    "dark": {
        "primary":    RGBColor(0x64, 0xB5, 0xF6),
        "secondary":  RGBColor(0x90, 0xCA, 0xF9),
        "accent":     RGBColor(0x00, 0xE6, 0x76),
        "bg":         RGBColor(0x1E, 0x1E, 0x2E),
        "bg_dark":    RGBColor(0x12, 0x12, 0x1E),
        "text":       RGBColor(0xE0, 0xE0, 0xE0),
        "text_light": RGBColor(0xFF, 0xFF, 0xFF),
        "text_muted": RGBColor(0x9E, 0x9E, 0x9E),
        "card_bg":    RGBColor(0x2A, 0x2A, 0x3E),
        "gradient_start": RGBColor(0x1E, 0x1E, 0x2E),
        "gradient_end":   RGBColor(0x12, 0x12, 0x1E),
        "decorations": ["stars", "grid", "glow"],
        "deco_color1": RGBColor(0x64, 0xB5, 0xF6),
        "deco_color2": RGBColor(0x00, 0xE6, 0x76),
        "deco_opacity": 0.20,
    },
    "green": {
        "primary":    RGBColor(0x2E, 0x7D, 0x32),
        "secondary":  RGBColor(0x1B, 0x5E, 0x20),
        "accent":     RGBColor(0x00, 0xC8, 0x53),
        "bg":         RGBColor(0xFF, 0xFF, 0xFF),
        "bg_dark":    RGBColor(0x1B, 0x5E, 0x20),
        "text":       RGBColor(0x21, 0x21, 0x21),
        "text_light": RGBColor(0xFF, 0xFF, 0xFF),
        "text_muted": RGBColor(0x75, 0x75, 0x75),
        "card_bg":    RGBColor(0xE8, 0xF5, 0xE9),
        "gradient_start": RGBColor(0x2E, 0x7D, 0x32),
        "gradient_end":   RGBColor(0x1B, 0x5E, 0x20),
        "decorations": ["leaves", "waves", "circles"],
        "deco_color1": RGBColor(0x2E, 0x7D, 0x32),
        "deco_color2": RGBColor(0x00, 0xC8, 0x53),
        "deco_opacity": 0.12,
    },
    "orange": {
        "primary":    RGBColor(0xE6, 0x51, 0x00),
        "secondary":  RGBColor(0xBF, 0x36, 0x0C),
        "accent":     RGBColor(0xFF, 0xAB, 0x00),
        "bg":         RGBColor(0xFF, 0xFF, 0xFF),
        "bg_dark":    RGBColor(0xBF, 0x36, 0x0C),
        "text":       RGBColor(0x21, 0x21, 0x21),
        "text_light": RGBColor(0xFF, 0xFF, 0xFF),
        "text_muted": RGBColor(0x75, 0x75, 0x75),
        "card_bg":    RGBColor(0xFF, 0xF3, 0xE0),
        "gradient_start": RGBColor(0xE6, 0x51, 0x00),
        "gradient_end":   RGBColor(0xBF, 0x36, 0x0C),
        "decorations": ["triangles", "waves", "dots"],
        "deco_color1": RGBColor(0xE6, 0x51, 0x00),
        "deco_color2": RGBColor(0xFF, 0xAB, 0x00),
        "deco_opacity": 0.15,
    },
    "purple": {
        "primary":    RGBColor(0x7B, 0x1F, 0xA2),
        "secondary":  RGBColor(0x4A, 0x14, 0x8C),
        "accent":     RGBColor(0xE0, 0x40, 0xFB),
        "bg":         RGBColor(0xFF, 0xFF, 0xFF),
        "bg_dark":    RGBColor(0x4A, 0x14, 0x8C),
        "text":       RGBColor(0x21, 0x21, 0x21),
        "text_light": RGBColor(0xFF, 0xFF, 0xFF),
        "text_muted": RGBColor(0x75, 0x75, 0x75),
        "card_bg":    RGBColor(0xF3, 0xE5, 0xF5),
        "gradient_start": RGBColor(0x7B, 0x1F, 0xA2),
        "gradient_end":   RGBColor(0x4A, 0x14, 0x8C),
        "decorations": ["circles", "stars", "glow"],
        "deco_color1": RGBColor(0x7B, 0x1F, 0xA2),
        "deco_color2": RGBColor(0xE0, 0x40, 0xFB),
        "deco_opacity": 0.15,
    },
    "red": {
        "primary":    RGBColor(0xC6, 0x28, 0x28),
        "secondary":  RGBColor(0xB7, 0x1C, 0x1C),
        "accent":     RGBColor(0xFF, 0x52, 0x52),
        "bg":         RGBColor(0xFF, 0xFF, 0xFF),
        "bg_dark":    RGBColor(0xB7, 0x1C, 0x1C),
        "text":       RGBColor(0x21, 0x21, 0x21),
        "text_light": RGBColor(0xFF, 0xFF, 0xFF),
        "text_muted": RGBColor(0x75, 0x75, 0x75),
        "card_bg":    RGBColor(0xFF, 0xEB, 0xEE),
        "gradient_start": RGBColor(0xC6, 0x28, 0x28),
        "gradient_end":   RGBColor(0xB7, 0x1C, 0x1C),
        "decorations": ["triangles", "waves", "circles"],
        "deco_color1": RGBColor(0xC6, 0x28, 0x28),
        "deco_color2": RGBColor(0xFF, 0x52, 0x52),
        "deco_opacity": 0.15,
    },
    "minimal": {
        "primary":    RGBColor(0x33, 0x33, 0x33),
        "secondary":  RGBColor(0x66, 0x66, 0x66),
        "accent":     RGBColor(0x00, 0x00, 0x00),
        "bg":         RGBColor(0xFF, 0xFF, 0xFF),
        "bg_dark":    RGBColor(0x33, 0x33, 0x33),
        "text":       RGBColor(0x33, 0x33, 0x33),
        "text_light": RGBColor(0xFF, 0xFF, 0xFF),
        "text_muted": RGBColor(0x99, 0x99, 0x99),
        "card_bg":    RGBColor(0xF5, 0xF5, 0xF5),
        "gradient_start": RGBColor(0x33, 0x33, 0x33),
        "gradient_end":   RGBColor(0x11, 0x11, 0x11),
        "decorations": ["lines", "dots"],
        "deco_color1": RGBColor(0x33, 0x33, 0x33),
        "deco_color2": RGBColor(0x99, 0x99, 0x99),
        "deco_opacity": 0.08,
    },
}

# ── Helper Functions ────────────────────────────────────────────────────

def _set_solid_bg(slide, color):
    """Set solid background color for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_shape(slide, shape_type, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    """Add a shape to the slide."""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def _add_textbox(slide, left, top, width, height, text, font_size=18,
                 font_color=None, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Microsoft YaHei"):
    """Add a text box to the slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    if font_color:
        p.font.color.rgb = font_color
    return txBox


def _add_accent_bar(slide, theme, left, top, width, height):
    """Add a decorative accent bar."""
    return _add_shape(slide, MSO_SHAPE.RECTANGLE, left, top, width, height, fill_color=theme["primary"])


def _add_icon_text(slide, icon, left, top, icon_size=36, text="", text_size=14,
                   icon_color=None, text_color=None, width=Inches(2)):
    """Add an icon (emoji) with text below it."""
    # Icon
    _add_textbox(slide, left, top, width, Inches(0.6), icon,
                 font_size=icon_size, font_color=icon_color,
                 alignment=PP_ALIGN.CENTER)
    # Text
    if text:
        _add_textbox(slide, left, top + Inches(0.6), width, Inches(0.8), text,
                     font_size=text_size, font_color=text_color,
                     alignment=PP_ALIGN.CENTER)


def _add_card(slide, theme, left, top, width, height, title="", bullets=None, icon=""):
    """Add a card-style element with optional icon, title, and bullets."""
    # Card background
    card = _add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height,
                      fill_color=theme["card_bg"])
    card.adjustments[0] = 0.05  # corner radius

    y_offset = top + Inches(0.2)

    # Icon
    if icon:
        _add_textbox(slide, left + Inches(0.2), y_offset, width - Inches(0.4), Inches(0.5),
                     icon, font_size=28, alignment=PP_ALIGN.LEFT)
        y_offset += Inches(0.5)

    # Title
    if title:
        _add_textbox(slide, left + Inches(0.2), y_offset, width - Inches(0.4), Inches(0.4),
                     title, font_size=16, font_color=theme["primary"], bold=True)
        y_offset += Inches(0.4)

    # Bullets
    if bullets:
        txBox = slide.shapes.add_textbox(left + Inches(0.2), y_offset,
                                         width - Inches(0.4), height - (y_offset - top) - Inches(0.1))
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"• {bullet}"
            p.font.size = Pt(12)
            p.font.color.rgb = theme["text"]
            p.space_after = Pt(4)


def _add_page_number(slide, theme, num, total):
    """Add page number to bottom-right corner."""
    _add_textbox(slide, Inches(8.8), Inches(5.2), Inches(1), Inches(0.3),
                 f"{num}/{total}", font_size=10, font_color=theme["text_muted"],
                 alignment=PP_ALIGN.RIGHT)


# ── Background Decoration System ────────────────────────────────────────

def _deco_circles(slide, theme, variant=0):
    """Add decorative circles to a slide. variant shifts positions."""
    c1 = theme.get("deco_color1", theme["primary"])
    c2 = theme.get("deco_color2", theme.get("accent", theme["primary"]))
    positions = [
        # variant 0: top-right + bottom-left
        [(Inches(7.5), Inches(-0.8), Inches(3.5), Inches(3.5), c1),
         (Inches(-1.2), Inches(3.8), Inches(2.5), Inches(2.5), c2)],
        # variant 1: top-left + bottom-right
        [(Inches(-0.8), Inches(-0.8), Inches(2.8), Inches(2.8), c2),
         (Inches(8.0), Inches(4.0), Inches(3.0), Inches(3.0), c1)],
        # variant 2: center-right + top-left small
        [(Inches(8.5), Inches(1.5), Inches(2.5), Inches(2.5), c1),
         (Inches(-0.5), Inches(-0.5), Inches(1.8), Inches(1.8), c2)],
        # variant 3: bottom-center + top-right small
        [(Inches(3.5), Inches(4.5), Inches(3.0), Inches(3.0), c2),
         (Inches(8.0), Inches(-0.5), Inches(2.0), Inches(2.0), c1)],
    ]
    v = variant % len(positions)
    for x, y, w, h, color in positions[v]:
        _add_shape(slide, MSO_SHAPE.OVAL, x, y, w, h, fill_color=color)


def _deco_dots(slide, theme, variant=0):
    """Add a grid of small decorative dots."""
    c1 = theme.get("deco_color1", theme["primary"])
    dot_size = Inches(0.12)
    # Place dots in a pattern based on variant
    start_x = Inches(8.0) if variant % 2 == 0 else Inches(0.3)
    start_y = Inches(0.5) if variant % 3 == 0 else Inches(4.0)
    for row in range(4):
        for col in range(4):
            x = start_x + col * Inches(0.3)
            y = start_y + row * Inches(0.3)
            _add_shape(slide, MSO_SHAPE.OVAL, x, y, dot_size, dot_size, fill_color=c1)


def _deco_triangles(slide, theme, variant=0):
    """Add decorative triangles."""
    c1 = theme.get("deco_color1", theme["primary"])
    c2 = theme.get("deco_color2", theme.get("accent", theme["primary"]))
    positions = [
        [(Inches(8.5), Inches(0), Inches(2), Inches(2), c1),
         (Inches(-0.5), Inches(4.2), Inches(1.5), Inches(1.5), c2)],
        [(Inches(0), Inches(0), Inches(1.8), Inches(1.8), c2),
         (Inches(8.8), Inches(4.5), Inches(1.8), Inches(1.8), c1)],
    ]
    v = variant % len(positions)
    for x, y, w, h, color in positions[v]:
        _add_shape(slide, MSO_SHAPE.ISOSCELES_TRIANGLE, x, y, w, h, fill_color=color)


def _deco_waves(slide, theme, variant=0):
    """Add wave-like horizontal bars at edges."""
    c1 = theme.get("deco_color1", theme["primary"])
    c2 = theme.get("deco_color2", theme.get("accent", theme["primary"]))
    if variant % 2 == 0:
        # Bottom waves
        _add_shape(slide, MSO_SHAPE.RECTANGLE,
                   Inches(0), Inches(5.2), Inches(10), Inches(0.08), fill_color=c1)
        _add_shape(slide, MSO_SHAPE.RECTANGLE,
                   Inches(0), Inches(5.35), Inches(7), Inches(0.05), fill_color=c2)
    else:
        # Top waves
        _add_shape(slide, MSO_SHAPE.RECTANGLE,
                   Inches(0), Inches(0), Inches(10), Inches(0.08), fill_color=c1)
        _add_shape(slide, MSO_SHAPE.RECTANGLE,
                   Inches(3), Inches(0.1), Inches(7), Inches(0.05), fill_color=c2)


def _deco_stars(slide, theme, variant=0):
    """Add star shapes as decorations."""
    c1 = theme.get("deco_color1", theme["primary"])
    c2 = theme.get("deco_color2", theme.get("accent", theme["primary"]))
    star_positions = [
        [(Inches(9.0), Inches(0.3), Inches(0.4), Inches(0.4), c1),
         (Inches(8.2), Inches(0.8), Inches(0.25), Inches(0.25), c2),
         (Inches(9.3), Inches(1.0), Inches(0.2), Inches(0.2), c1)],
        [(Inches(0.3), Inches(4.5), Inches(0.35), Inches(0.35), c2),
         (Inches(0.8), Inches(5.0), Inches(0.2), Inches(0.2), c1),
         (Inches(0.1), Inches(5.1), Inches(0.25), Inches(0.25), c2)],
    ]
    v = variant % len(star_positions)
    for x, y, w, h, color in star_positions[v]:
        _add_shape(slide, MSO_SHAPE.STAR_4_POINT, x, y, w, h, fill_color=color)


def _deco_grid(slide, theme, variant=0):
    """Add a subtle grid pattern using thin lines."""
    c1 = theme.get("deco_color1", theme["primary"])
    # Vertical lines
    start_x = Inches(8.5) if variant % 2 == 0 else Inches(0.2)
    for i in range(5):
        x = start_x + i * Inches(0.25)
        _add_shape(slide, MSO_SHAPE.RECTANGLE,
                   x, Inches(0.3), Inches(0.015), Inches(1.5), fill_color=c1)
    # Horizontal lines
    for i in range(4):
        y = Inches(0.3) + i * Inches(0.4)
        _add_shape(slide, MSO_SHAPE.RECTANGLE,
                   start_x, y, Inches(1.2), Inches(0.015), fill_color=c1)


def _deco_glow(slide, theme, variant=0):
    """Add glow-like large transparent circles."""
    c1 = theme.get("deco_color1", theme["primary"])
    c2 = theme.get("deco_color2", theme.get("accent", theme["primary"]))
    glow_positions = [
        [(Inches(7.0), Inches(-1.5), Inches(5.0), Inches(5.0), c1),
         (Inches(-2.0), Inches(2.5), Inches(4.5), Inches(4.5), c2)],
        [(Inches(-1.5), Inches(-1.5), Inches(4.5), Inches(4.5), c1),
         (Inches(6.5), Inches(3.0), Inches(5.0), Inches(5.0), c2)],
    ]
    v = variant % len(glow_positions)
    for x, y, w, h, color in glow_positions[v]:
        _add_shape(slide, MSO_SHAPE.OVAL, x, y, w, h, fill_color=color)


def _deco_lines(slide, theme, variant=0):
    """Add minimal decorative lines."""
    c1 = theme.get("deco_color1", theme["primary"])
    c2 = theme.get("deco_color2", theme.get("accent", theme["primary"]))
    if variant % 2 == 0:
        # Right side vertical line
        _add_shape(slide, MSO_SHAPE.RECTANGLE,
                   Inches(9.5), Inches(0.5), Inches(0.03), Inches(4.5), fill_color=c1)
        _add_shape(slide, MSO_SHAPE.RECTANGLE,
                   Inches(9.6), Inches(1.0), Inches(0.03), Inches(3.5), fill_color=c2)
    else:
        # Left side vertical line
        _add_shape(slide, MSO_SHAPE.RECTANGLE,
                   Inches(0.3), Inches(0.5), Inches(0.03), Inches(4.5), fill_color=c1)
        _add_shape(slide, MSO_SHAPE.RECTANGLE,
                   Inches(0.4), Inches(1.0), Inches(0.03), Inches(3.5), fill_color=c2)


def _deco_leaves(slide, theme, variant=0):
    """Add leaf-like oval shapes as organic decorations."""
    c1 = theme.get("deco_color1", theme["primary"])
    c2 = theme.get("deco_color2", theme.get("accent", theme["primary"]))
    leaf_positions = [
        [(Inches(8.5), Inches(0.2), Inches(1.8), Inches(0.8), c1),
         (Inches(-0.3), Inches(4.5), Inches(1.5), Inches(0.7), c2)],
        [(Inches(0.2), Inches(0.2), Inches(1.5), Inches(0.7), c2),
         (Inches(8.8), Inches(4.8), Inches(1.8), Inches(0.8), c1)],
    ]
    v = variant % len(leaf_positions)
    for x, y, w, h, color in leaf_positions[v]:
        _add_shape(slide, MSO_SHAPE.OVAL, x, y, w, h, fill_color=color)


# Map decoration names to functions
DECO_FUNCS = {
    "circles":   _deco_circles,
    "dots":      _deco_dots,
    "triangles": _deco_triangles,
    "waves":     _deco_waves,
    "stars":     _deco_stars,
    "grid":      _deco_grid,
    "glow":      _deco_glow,
    "lines":     _deco_lines,
    "leaves":    _deco_leaves,
}


def _add_decorations(slide, theme, slide_type="content", slide_index=0):
    """
    Add theme-specific decorations to a slide.
    
    Different slide types get different decoration styles:
    - title/section/end: heavy decorations (glow, circles, stars)
    - content/data slides: lighter decorations (dots, waves, lines)
    
    slide_index is used as variant to vary positions across slides.
    """
    deco_list = theme.get("decorations", [])
    if not deco_list:
        return
    
    # Choose decorations based on slide type
    if slide_type in ("title", "section", "end"):
        # Use the first and last decoration types for emphasis
        chosen = [deco_list[0]]
        if len(deco_list) > 2:
            chosen.append(deco_list[2])
    else:
        # Use lighter decorations for content slides
        chosen = deco_list[1:2] if len(deco_list) > 1 else deco_list[:1]
    
    for deco_name in chosen:
        func = DECO_FUNCS.get(deco_name)
        if func:
            func(slide, theme, variant=slide_index)


# ── Slide Builders ──────────────────────────────────────────────────────

def build_title_slide(prs, data, theme, slide_num, total):
    """Title slide with gradient background and decorative elements."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _set_solid_bg(slide, theme["bg_dark"])

    # Add theme-specific decorations (heavy for title)
    _add_decorations(slide, theme, slide_type="title", slide_index=slide_num)

    # Decorative accent bar at top
    _add_shape(slide, MSO_SHAPE.RECTANGLE,
               Inches(0), Inches(0), Inches(10), Inches(0.08),
               fill_color=theme["accent"])

    # Decorative circle (top-right)
    _add_shape(slide, MSO_SHAPE.OVAL,
               Inches(8), Inches(-0.5), Inches(2.5), Inches(2.5),
               fill_color=theme["primary"])

    # Decorative circle (bottom-left)
    _add_shape(slide, MSO_SHAPE.OVAL,
               Inches(-0.8), Inches(4), Inches(2), Inches(2),
               fill_color=theme["secondary"])

    # Title
    _add_textbox(slide, Inches(1), Inches(1.5), Inches(8), Inches(1.2),
                 data.get("title", ""), font_size=40, font_color=theme["text_light"],
                 bold=True, alignment=PP_ALIGN.LEFT)

    # Subtitle
    subtitle = data.get("subtitle", "")
    if subtitle:
        _add_textbox(slide, Inches(1), Inches(2.8), Inches(7), Inches(0.6),
                     subtitle, font_size=22, font_color=theme["text_light"],
                     alignment=PP_ALIGN.LEFT)

    # Accent line under title
    _add_shape(slide, MSO_SHAPE.RECTANGLE,
               Inches(1), Inches(2.6), Inches(2), Inches(0.06),
               fill_color=theme["accent"])

    # Author / date
    author = data.get("author", "")
    date = data.get("date", "")
    footer_text = f"{author}  |  {date}" if author and date else (author or date)
    if footer_text:
        _add_textbox(slide, Inches(1), Inches(4.5), Inches(6), Inches(0.4),
                     footer_text, font_size=14, font_color=theme["text_light"],
                     alignment=PP_ALIGN.LEFT)

    _add_page_number(slide, theme, slide_num, total)
    return slide


def build_section_slide(prs, data, theme, slide_num, total):
    """Section header with colored background."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_solid_bg(slide, theme["bg_dark"])

    # Add theme-specific decorations
    _add_decorations(slide, theme, slide_type="section", slide_index=slide_num)

    # Decorative accent bar
    _add_shape(slide, MSO_SHAPE.RECTANGLE,
               Inches(1), Inches(2.3), Inches(1.5), Inches(0.06),
               fill_color=theme["accent"])

    # Section icon (optional)
    icon = data.get("icon", "")
    if icon:
        _add_textbox(slide, Inches(1), Inches(1.2), Inches(2), Inches(0.8),
                     icon, font_size=48, font_color=theme["text_light"],
                     alignment=PP_ALIGN.LEFT)

    # Section title
    _add_textbox(slide, Inches(1), Inches(2.6), Inches(8), Inches(1),
                 data.get("title", ""), font_size=36, font_color=theme["text_light"],
                 bold=True, alignment=PP_ALIGN.LEFT)

    # Section subtitle (optional)
    subtitle = data.get("subtitle", "")
    if subtitle:
        _add_textbox(slide, Inches(1), Inches(3.6), Inches(7), Inches(0.5),
                     subtitle, font_size=18, font_color=theme["text_light"],
                     alignment=PP_ALIGN.LEFT)

    _add_page_number(slide, theme, slide_num, total)
    return slide


def build_content_slide(prs, data, theme, slide_num, total):
    """Content slide with accent bar, icon, and bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_solid_bg(slide, theme["bg"])

    # Add theme-specific decorations (lighter for content)
    _add_decorations(slide, theme, slide_type="content", slide_index=slide_num)

    # Top accent bar
    _add_shape(slide, MSO_SHAPE.RECTANGLE,
               Inches(0), Inches(0), Inches(10), Inches(0.06),
               fill_color=theme["primary"])

    # Left accent bar
    _add_shape(slide, MSO_SHAPE.RECTANGLE,
               Inches(0), Inches(0), Inches(0.06), Inches(5.625),
               fill_color=theme["primary"])

    # Slide icon (optional)
    icon = data.get("icon", "")
    title_left = Inches(0.8)
    if icon:
        _add_textbox(slide, Inches(0.8), Inches(0.3), Inches(0.8), Inches(0.6),
                     icon, font_size=32, alignment=PP_ALIGN.LEFT)
        title_left = Inches(1.5)

    # Title
    _add_textbox(slide, title_left, Inches(0.35), Inches(8), Inches(0.6),
                 data.get("title", ""), font_size=28, font_color=theme["primary"],
                 bold=True, alignment=PP_ALIGN.LEFT)

    # Title underline
    _add_shape(slide, MSO_SHAPE.RECTANGLE,
               Inches(0.8), Inches(1.05), Inches(8.4), Inches(0.02),
               fill_color=theme["card_bg"])

    # Bullet points
    bullets = data.get("bullets", [])
    if bullets:
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.4), Inches(3.8))
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            # Support bullet icons
            bullet_icon = data.get("bullet_icon", "▸")
            p.text = f"  {bullet_icon}  {bullet}"
            p.font.size = Pt(18)
            p.font.color.rgb = theme["text"]
            p.font.name = "Microsoft YaHei"
            p.space_after = Pt(12)
            p.space_before = Pt(4)

    # Footer note (optional)
    note = data.get("note", "")
    if note:
        _add_textbox(slide, Inches(0.8), Inches(5.0), Inches(8), Inches(0.4),
                     note, font_size=11, font_color=theme["text_muted"],
                     alignment=PP_ALIGN.LEFT)

    _add_page_number(slide, theme, slide_num, total)
    return slide


def build_two_column_slide(prs, data, theme, slide_num, total):
    """Two-column layout with cards."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_solid_bg(slide, theme["bg"])

    # Add theme-specific decorations
    _add_decorations(slide, theme, slide_type="content", slide_index=slide_num)

    # Top accent bar
    _add_shape(slide, MSO_SHAPE.RECTANGLE,
               Inches(0), Inches(0), Inches(10), Inches(0.06),
               fill_color=theme["primary"])

    # Title
    _add_textbox(slide, Inches(0.8), Inches(0.35), Inches(8), Inches(0.6),
                 data.get("title", ""), font_size=28, font_color=theme["primary"],
                 bold=True, alignment=PP_ALIGN.LEFT)

    # Title underline
    _add_shape(slide, MSO_SHAPE.RECTANGLE,
               Inches(0.8), Inches(1.05), Inches(8.4), Inches(0.02),
               fill_color=theme["card_bg"])

    left = data.get("left", {})
    right = data.get("right", {})

    # Left card
    _add_card(slide, theme,
              Inches(0.5), Inches(1.3), Inches(4.3), Inches(3.8),
              title=left.get("title", ""),
              bullets=left.get("bullets", []),
              icon=left.get("icon", ""))

    # Right card
    _add_card(slide, theme,
              Inches(5.2), Inches(1.3), Inches(4.3), Inches(3.8),
              title=right.get("title", ""),
              bullets=right.get("bullets", []),
              icon=right.get("icon", ""))

    _add_page_number(slide, theme, slide_num, total)
    return slide


def build_icon_grid_slide(prs, data, theme, slide_num, total):
    """Grid of icon + title + description cards."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_solid_bg(slide, theme["bg"])

    # Add theme-specific decorations
    _add_decorations(slide, theme, slide_type="content", slide_index=slide_num)

    # Top accent bar
    _add_shape(slide, MSO_SHAPE.RECTANGLE,
               Inches(0), Inches(0), Inches(10), Inches(0.06),
               fill_color=theme["primary"])

    # Title
    _add_textbox(slide, Inches(0.8), Inches(0.35), Inches(8), Inches(0.6),
                 data.get("title", ""), font_size=28, font_color=theme["primary"],
                 bold=True, alignment=PP_ALIGN.LEFT)

    # Title underline
    _add_shape(slide, MSO_SHAPE.RECTANGLE,
               Inches(0.8), Inches(1.05), Inches(8.4), Inches(0.02),
               fill_color=theme["card_bg"])

    items = data.get("items", [])
    cols = data.get("columns", min(len(items), 3))
    rows = (len(items) + cols - 1) // cols

    card_w = Inches(8.4 / cols - 0.2)
    card_h = Inches(min(3.5 / rows, 1.8))
    start_x = Inches(0.5)
    start_y = Inches(1.3)
    gap_x = Inches(8.4 / cols)
    gap_y = Inches(0.15) + card_h

    for idx, item in enumerate(items):
        col = idx % cols
        row = idx // cols
        x = start_x + col * gap_x
        y = start_y + row * gap_y

        _add_card(slide, theme, x, y, card_w, card_h,
                  title=item.get("title", ""),
                  bullets=item.get("bullets", []),
                  icon=item.get("icon", ""))

    _add_page_number(slide, theme, slide_num, total)
    return slide


def build_stats_slide(prs, data, theme, slide_num, total):
    """Statistics / KPI slide with large numbers."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_solid_bg(slide, theme["bg"])

    # Add theme-specific decorations
    _add_decorations(slide, theme, slide_type="content", slide_index=slide_num)

    # Top accent bar
    _add_shape(slide, MSO_SHAPE.RECTANGLE,
               Inches(0), Inches(0), Inches(10), Inches(0.06),
               fill_color=theme["primary"])

    # Title
    _add_textbox(slide, Inches(0.8), Inches(0.35), Inches(8), Inches(0.6),
                 data.get("title", ""), font_size=28, font_color=theme["primary"],
                 bold=True, alignment=PP_ALIGN.LEFT)

    # Title underline
    _add_shape(slide, MSO_SHAPE.RECTANGLE,
               Inches(0.8), Inches(1.05), Inches(8.4), Inches(0.02),
               fill_color=theme["card_bg"])

    stats = data.get("stats", [])
    cols = min(len(stats), 4)
    card_w = Inches(8.4 / cols - 0.2)
    gap_x = Inches(8.4 / cols)
    start_x = Inches(0.5)
    start_y = Inches(1.5)

    for idx, stat in enumerate(stats):
        col = idx % cols
        x = start_x + col * gap_x

        # Card background
        _add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                   x, start_y, card_w, Inches(3.2),
                   fill_color=theme["card_bg"])

        # Icon
        icon = stat.get("icon", "")
        if icon:
            _add_textbox(slide, x, start_y + Inches(0.2), card_w, Inches(0.5),
                         icon, font_size=32, alignment=PP_ALIGN.CENTER)

        # Number
        _add_textbox(slide, x, start_y + Inches(0.8), card_w, Inches(0.8),
                     stat.get("value", ""), font_size=36, font_color=theme["primary"],
                     bold=True, alignment=PP_ALIGN.CENTER)

        # Label
        _add_textbox(slide, x, start_y + Inches(1.6), card_w, Inches(0.4),
                     stat.get("label", ""), font_size=14, font_color=theme["text"],
                     alignment=PP_ALIGN.CENTER)

        # Description
        desc = stat.get("description", "")
        if desc:
            _add_textbox(slide, x + Inches(0.1), start_y + Inches(2.1),
                         card_w - Inches(0.2), Inches(0.8),
                         desc, font_size=11, font_color=theme["text_muted"],
                         alignment=PP_ALIGN.CENTER)

    _add_page_number(slide, theme, slide_num, total)
    return slide


def build_chart_slide(prs, data, theme, slide_num, total):
    """Slide with a chart (bar, pie, line)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_solid_bg(slide, theme["bg"])

    # Add theme-specific decorations
    _add_decorations(slide, theme, slide_type="content", slide_index=slide_num)

    # Top accent bar
    _add_shape(slide, MSO_SHAPE.RECTANGLE,
               Inches(0), Inches(0), Inches(10), Inches(0.06),
               fill_color=theme["primary"])

    # Title
    _add_textbox(slide, Inches(0.8), Inches(0.35), Inches(8), Inches(0.6),
                 data.get("title", ""), font_size=28, font_color=theme["primary"],
                 bold=True, alignment=PP_ALIGN.LEFT)

    chart_type = data.get("chart_type", "bar")
    chart_data = CategoryChartData()
    chart_data.categories = data.get("categories", [])
    for series in data.get("series", []):
        chart_data.add_series(series.get("name", ""), series.get("values", []))

    type_map = {
        "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "line": XL_CHART_TYPE.LINE,
        "pie": XL_CHART_TYPE.PIE,
        "area": XL_CHART_TYPE.AREA,
    }
    xl_type = type_map.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)

    chart_frame = slide.shapes.add_chart(
        xl_type, Inches(1), Inches(1.3), Inches(8), Inches(4),
        chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = len(data.get("series", [])) > 1
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False

    # Apply theme colors to series
    colors = [theme["primary"], theme["accent"], theme["secondary"],
              RGBColor(0xFF, 0x98, 0x00), RGBColor(0x4C, 0xAF, 0x50)]
    for i, series in enumerate(chart.series):
        color = colors[i % len(colors)]
        if chart_type == "pie":
            # Color each point differently for pie
            for j in range(len(data.get("categories", []))):
                point = series.points[j]
                point.format.fill.solid()
                point.format.fill.fore_color.rgb = colors[j % len(colors)]
        else:
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = color

    _add_page_number(slide, theme, slide_num, total)
    return slide


def build_table_slide(prs, data, theme, slide_num, total):
    """Slide with a styled data table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_solid_bg(slide, theme["bg"])

    # Add theme-specific decorations
    _add_decorations(slide, theme, slide_type="content", slide_index=slide_num)

    # Top accent bar
    _add_shape(slide, MSO_SHAPE.RECTANGLE,
               Inches(0), Inches(0), Inches(10), Inches(0.06),
               fill_color=theme["primary"])

    # Title
    _add_textbox(slide, Inches(0.8), Inches(0.35), Inches(8), Inches(0.6),
                 data.get("title", ""), font_size=28, font_color=theme["primary"],
                 bold=True, alignment=PP_ALIGN.LEFT)

    headers = data.get("headers", [])
    rows_data = data.get("rows", [])
    n_rows = len(rows_data) + 1  # +1 for header
    n_cols = len(headers)

    table_width = Inches(8.4)
    table_height = Inches(min(0.4 * n_rows, 4.0))
    left = Inches(0.8)
    top = Inches(1.3)

    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, table_width, table_height)
    table = table_shape.table

    # Header row
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = header
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(13)
            paragraph.font.bold = True
            paragraph.font.color.rgb = theme["text_light"]
            paragraph.font.name = "Microsoft YaHei"
        cell.fill.solid()
        cell.fill.fore_color.rgb = theme["primary"]

    # Data rows
    for i, row in enumerate(rows_data):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(12)
                paragraph.font.color.rgb = theme["text"]
                paragraph.font.name = "Microsoft YaHei"
            cell.fill.solid()
            cell.fill.fore_color.rgb = theme["bg"] if i % 2 == 0 else theme["card_bg"]

    _add_page_number(slide, theme, slide_num, total)
    return slide


def build_image_slide(prs, data, theme, slide_num, total):
    """Slide with an image."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_solid_bg(slide, theme["bg"])

    # Top accent bar
    _add_shape(slide, MSO_SHAPE.RECTANGLE,
               Inches(0), Inches(0), Inches(10), Inches(0.06),
               fill_color=theme["primary"])

    # Title
    _add_textbox(slide, Inches(0.8), Inches(0.35), Inches(8), Inches(0.6),
                 data.get("title", ""), font_size=28, font_color=theme["primary"],
                 bold=True, alignment=PP_ALIGN.LEFT)

    image_path = data.get("image_path", "")
    if image_path and Path(image_path).exists():
        slide.shapes.add_picture(image_path, Inches(1.5), Inches(1.3), width=Inches(7))
    else:
        _add_textbox(slide, Inches(2), Inches(2.5), Inches(6), Inches(1),
                     f"[Image: {image_path or 'not specified'}]",
                     font_size=16, font_color=theme["text_muted"],
                     alignment=PP_ALIGN.CENTER)

    caption = data.get("caption", "")
    if caption:
        _add_textbox(slide, Inches(1), Inches(4.8), Inches(8), Inches(0.4),
                     caption, font_size=12, font_color=theme["text_muted"],
                     alignment=PP_ALIGN.CENTER)

    _add_page_number(slide, theme, slide_num, total)
    return slide


def build_timeline_slide(prs, data, theme, slide_num, total):
    """Timeline slide with icon milestones."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_solid_bg(slide, theme["bg"])

    # Top accent bar
    _add_shape(slide, MSO_SHAPE.RECTANGLE,
               Inches(0), Inches(0), Inches(10), Inches(0.06),
               fill_color=theme["primary"])

    # Title
    _add_textbox(slide, Inches(0.8), Inches(0.35), Inches(8), Inches(0.6),
                 data.get("title", ""), font_size=28, font_color=theme["primary"],
                 bold=True, alignment=PP_ALIGN.LEFT)

    # Title underline
    _add_shape(slide, MSO_SHAPE.RECTANGLE,
               Inches(0.8), Inches(1.05), Inches(8.4), Inches(0.02),
               fill_color=theme["card_bg"])

    events = data.get("events", [])
    if not events:
        return slide

    # Timeline line
    line_y = Inches(2.8)
    _add_shape(slide, MSO_SHAPE.RECTANGLE,
               Inches(0.8), line_y, Inches(8.4), Inches(0.04),
               fill_color=theme["primary"])

    # Events
    n = len(events)
    spacing = Inches(8.4 / n)
    for i, event in enumerate(events):
        x = Inches(0.8) + i * spacing + spacing / 2 - Inches(0.3)

        # Circle node
        _add_shape(slide, MSO_SHAPE.OVAL,
                   x, line_y - Inches(0.15), Inches(0.3), Inches(0.3),
                   fill_color=theme["accent"])

        # Icon above
        icon = event.get("icon", "")
        if icon:
            _add_textbox(slide, x - Inches(0.2), line_y - Inches(1.0), Inches(0.7), Inches(0.6),
                         icon, font_size=32, alignment=PP_ALIGN.CENTER)

        # Year/date above
        date = event.get("date", "")
        if date:
            _add_textbox(slide, x - Inches(0.3), line_y - Inches(1.5), Inches(0.9), Inches(0.4),
                         date, font_size=14, font_color=theme["primary"],
                         bold=True, alignment=PP_ALIGN.CENTER)

        # Description below
        desc = event.get("description", "")
        if desc:
            _add_textbox(slide, x - Inches(0.4), line_y + Inches(0.3), Inches(1.1), Inches(1.2),
                         desc, font_size=11, font_color=theme["text"],
                         alignment=PP_ALIGN.CENTER)

    _add_page_number(slide, theme, slide_num, total)
    return slide


def build_process_slide(prs, data, theme, slide_num, total):
    """Process flow slide with numbered/icon steps."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_solid_bg(slide, theme["bg"])

    # Top accent bar
    _add_shape(slide, MSO_SHAPE.RECTANGLE,
               Inches(0), Inches(0), Inches(10), Inches(0.06),
               fill_color=theme["primary"])

    # Title
    _add_textbox(slide, Inches(0.8), Inches(0.35), Inches(8), Inches(0.6),
                 data.get("title", ""), font_size=28, font_color=theme["primary"],
                 bold=True, alignment=PP_ALIGN.LEFT)

    # Title underline
    _add_shape(slide, MSO_SHAPE.RECTANGLE,
               Inches(0.8), Inches(1.05), Inches(8.4), Inches(0.02),
               fill_color=theme["card_bg"])

    steps = data.get("steps", [])
    if not steps:
        return slide

    n = len(steps)
    card_w = Inches(8.4 / n - 0.15)
    gap = Inches(8.4 / n)
    start_x = Inches(0.5)
    start_y = Inches(1.5)

    for i, step in enumerate(steps):
        x = start_x + i * gap

        # Card
        _add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                   x, start_y, card_w, Inches(3.5),
                   fill_color=theme["card_bg"])

        # Step number circle
        _add_shape(slide, MSO_SHAPE.OVAL,
                   x + card_w / 2 - Inches(0.25), start_y + Inches(0.2),
                   Inches(0.5), Inches(0.5),
                   fill_color=theme["primary"])

        _add_textbox(slide, x + card_w / 2 - Inches(0.25), start_y + Inches(0.25),
                     Inches(0.5), Inches(0.4),
                     str(i + 1), font_size=20, font_color=theme["text_light"],
                     bold=True, alignment=PP_ALIGN.CENTER)

        # Icon
        icon = step.get("icon", "")
        if icon:
            _add_textbox(slide, x, start_y + Inches(0.9), card_w, Inches(0.5),
                         icon, font_size=32, alignment=PP_ALIGN.CENTER)

        # Title
        title = step.get("title", "")
        if title:
            _add_textbox(slide, x + Inches(0.1), start_y + Inches(1.5),
                         card_w - Inches(0.2), Inches(0.4),
                         title, font_size=14, font_color=theme["primary"],
                         bold=True, alignment=PP_ALIGN.CENTER)

        # Description
        desc = step.get("description", "")
        if desc:
            _add_textbox(slide, x + Inches(0.1), start_y + Inches(2.0),
                         card_w - Inches(0.2), Inches(1.2),
                         desc, font_size=11, font_color=theme["text"],
                         alignment=PP_ALIGN.CENTER)

        # Arrow between steps
        if i < n - 1:
            arrow_x = x + card_w + Inches(0.05)
            _add_textbox(slide, arrow_x, start_y + Inches(1.5),
                         Inches(0.2), Inches(0.4),
                         "→", font_size=24, font_color=theme["primary"],
                         bold=True, alignment=PP_ALIGN.CENTER)

    _add_page_number(slide, theme, slide_num, total)
    return slide


def build_comparison_slide(prs, data, theme, slide_num, total):
    """Comparison slide with icon-based feature cards."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_solid_bg(slide, theme["bg"])

    # Top accent bar
    _add_shape(slide, MSO_SHAPE.RECTANGLE,
               Inches(0), Inches(0), Inches(10), Inches(0.06),
               fill_color=theme["primary"])

    # Title
    _add_textbox(slide, Inches(0.8), Inches(0.35), Inches(8), Inches(0.6),
                 data.get("title", ""), font_size=28, font_color=theme["primary"],
                 bold=True, alignment=PP_ALIGN.LEFT)

    # Title underline
    _add_shape(slide, MSO_SHAPE.RECTANGLE,
               Inches(0.8), Inches(1.05), Inches(8.4), Inches(0.02),
               fill_color=theme["card_bg"])

    items = data.get("items", [])
    if not items:
        return slide

    n = len(items)
    cols = min(n, 3)
    rows = (n + cols - 1) // cols
    card_w = Inches(8.4 / cols - 0.2)
    card_h = Inches(min(3.5 / rows, 2.2))
    gap_x = Inches(8.4 / cols)
    gap_y = Inches(0.15) + card_h
    start_x = Inches(0.5)
    start_y = Inches(1.3)

    for idx, item in enumerate(items):
        col = idx % cols
        row = idx // cols
        x = start_x + col * gap_x
        y = start_y + row * gap_y

        # Card background
        _add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                   x, y, card_w, card_h,
                   fill_color=theme["card_bg"])

        # Icon
        icon = item.get("icon", "")
        if icon:
            _add_textbox(slide, x, y + Inches(0.2), card_w, Inches(0.5),
                         icon, font_size=36, alignment=PP_ALIGN.CENTER)

        # Title
        title = item.get("title", "")
        if title:
            _add_textbox(slide, x + Inches(0.1), y + Inches(0.8),
                         card_w - Inches(0.2), Inches(0.4),
                         title, font_size=16, font_color=theme["primary"],
                         bold=True, alignment=PP_ALIGN.CENTER)

        # Value/metric
        value = item.get("value", "")
        if value:
            _add_textbox(slide, x + Inches(0.1), y + Inches(1.2),
                         card_w - Inches(0.2), Inches(0.5),
                         value, font_size=24, font_color=theme["accent"],
                         bold=True, alignment=PP_ALIGN.CENTER)

        # Description
        desc = item.get("description", "")
        if desc:
            _add_textbox(slide, x + Inches(0.1), y + Inches(1.7),
                         card_w - Inches(0.2), Inches(0.8),
                         desc, font_size=11, font_color=theme["text"],
                         alignment=PP_ALIGN.CENTER)

    _add_page_number(slide, theme, slide_num, total)
    return slide


def build_end_slide(prs, data, theme, slide_num, total):
    """Thank you / ending slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_solid_bg(slide, theme["bg_dark"])

    # Decorative accent bar at bottom
    _add_shape(slide, MSO_SHAPE.RECTANGLE,
               Inches(0), Inches(5.545), Inches(10), Inches(0.08),
               fill_color=theme["accent"])

    # Decorative circle
    _add_shape(slide, MSO_SHAPE.OVAL,
               Inches(-1), Inches(-1), Inches(3), Inches(3),
               fill_color=theme["primary"])

    _add_shape(slide, MSO_SHAPE.OVAL,
               Inches(8.5), Inches(3.5), Inches(2.5), Inches(2.5),
               fill_color=theme["secondary"])

    # Thank you text
    title = data.get("title", "Thank You")
    _add_textbox(slide, Inches(1), Inches(1.5), Inches(8), Inches(1.2),
                 title, font_size=44, font_color=theme["text_light"],
                 bold=True, alignment=PP_ALIGN.CENTER)

    # Accent line
    _add_shape(slide, MSO_SHAPE.RECTANGLE,
               Inches(4), Inches(2.8), Inches(2), Inches(0.06),
               fill_color=theme["accent"])

    # Contact / subtitle
    subtitle = data.get("subtitle", "")
    if subtitle:
        _add_textbox(slide, Inches(1), Inches(3.1), Inches(8), Inches(0.6),
                     subtitle, font_size=18, font_color=theme["text_light"],
                     alignment=PP_ALIGN.CENTER)

    # Contact info
    contact = data.get("contact", "")
    if contact:
        _add_textbox(slide, Inches(1), Inches(4.0), Inches(8), Inches(0.5),
                     contact, font_size=14, font_color=theme["text_light"],
                     alignment=PP_ALIGN.CENTER)

    _add_page_number(slide, theme, slide_num, total)
    return slide


# ── Slide Type Router ───────────────────────────────────────────────────

BUILDERS = {
    "title":      build_title_slide,
    "section":    build_section_slide,
    "content":    build_content_slide,
    "two_column": build_two_column_slide,
    "icon_grid":  build_icon_grid_slide,
    "stats":      build_stats_slide,
    "chart":      build_chart_slide,
    "table":      build_table_slide,
    "image":      build_image_slide,
    "timeline":   build_timeline_slide,
    "process":    build_process_slide,
    "comparison": build_comparison_slide,
    "end":        build_end_slide,
}


def generate_presentation(data, output_path, theme_name="blue"):
    """Generate a PowerPoint presentation from structured data."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    theme = THEMES.get(theme_name, THEMES["blue"])

    # Allow per-presentation theme override
    if "theme" in data:
        theme = THEMES.get(data["theme"], theme)

    slides_data = []

    # Auto-add title slide if top-level title exists
    if "title" in data:
        slides_data.append({
            "type": "title",
            "title": data["title"],
            "subtitle": data.get("subtitle", ""),
            "author": data.get("author", ""),
            "date": data.get("date", ""),
        })

    slides_data.extend(data.get("slides", []))

    # Auto-add end slide if last slide is not type "end"
    if slides_data and slides_data[-1].get("type") != "end":
        slides_data.append({
            "type": "end",
            "title": data.get("end_title", "Thank You"),
            "subtitle": data.get("end_subtitle", ""),
            "contact": data.get("end_contact", ""),
        })

    total = len(slides_data)

    for i, slide_data in enumerate(slides_data):
        slide_type = slide_data.get("type", "content")
        builder = BUILDERS.get(slide_type, build_content_slide)
        builder(prs, slide_data, theme, i + 1, total)

    prs.save(output_path)
    print(f"✅ Presentation saved: {output_path}")
    print(f"   Theme: {theme_name}")
    print(f"   Total slides: {len(prs.slides)}")


def main():
    if len(sys.argv) < 3:
        print("Usage: python generate_ppt.py <input.json> <output.pptx> [--theme <theme>]")
        print(f"Available themes: {', '.join(THEMES.keys())}")
        sys.exit(1)

    input_path = sys.argv[1]
    output_path = sys.argv[2]
    theme_name = "blue"

    if "--theme" in sys.argv:
        idx = sys.argv.index("--theme")
        if idx + 1 < len(sys.argv):
            theme_name = sys.argv[idx + 1]

    try:
        with open(input_path, 'r', encoding='utf-8') as f:
            data = json.load(f)
    except FileNotFoundError:
        print(f"Error: Input file not found: {input_path}")
        sys.exit(1)
    except json.JSONDecodeError as e:
        print(f"Error: Invalid JSON: {e}")
        sys.exit(1)

    generate_presentation(data, output_path, theme_name)


if __name__ == "__main__":
    main()
